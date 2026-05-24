from __future__ import annotations

import sys
from datetime import datetime
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
SHOWCASE_DIR = ROOT / "state" / "showcase"
PPTX_DEPS = SHOWCASE_DIR / "pptx_deps"
OUTPUT = SHOWCASE_DIR / "RRKAL_Showcase_Guide.zh-TW.pptx"

if PPTX_DEPS.exists():
    sys.path.insert(0, str(PPTX_DEPS))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402


FONT = "Microsoft JhengHei"
BG = RGBColor(12, 18, 26)
PANEL = RGBColor(22, 31, 43)
PANEL_2 = RGBColor(31, 45, 61)
ACCENT = RGBColor(56, 189, 248)
GREEN = RGBColor(74, 222, 128)
YELLOW = RGBColor(250, 204, 21)
RED = RGBColor(248, 113, 113)
TEXT = RGBColor(235, 245, 255)
MUTED = RGBColor(156, 178, 201)
INK = RGBColor(9, 14, 21)


def set_background(slide, color=BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, text: str, x: float, y: float, w: float, h: float, *, size: int = 24, color=TEXT, bold: bool = False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, *, size: int = 20, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = color
        paragraph.space_after = Pt(8)
    return box


def add_panel(slide, x: float, y: float, w: float, h: float, *, fill=PANEL, line=None, radius: bool = True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or RGBColor(48, 65, 84)
    shape.line.width = Pt(1)
    return shape


def add_tag(slide, text: str, x: float, y: float, *, color=ACCENT):
    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.65), Inches(0.34))
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.color.rgb = color
    frame = tag.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = INK
    return tag


def add_progress(slide, x: float, y: float, w: float, pct: float, label: str, *, color=GREEN) -> None:
    add_panel(slide, x, y, w, 0.18, fill=RGBColor(45, 56, 72), line=RGBColor(45, 56, 72))
    add_panel(slide, x, y, max(0.01, w * pct / 100), 0.18, fill=color, line=color)
    add_text(slide, f"{label} {pct:.0f}%", x, y - 0.35, w, 0.28, size=12, color=MUTED)


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_tag(slide, "SHOWCASE MODE", 0.65, 0.45)
    add_text(slide, "RuRuKa Asset Launcher", 0.65, 1.05, 8.0, 0.7, size=42, bold=True)
    add_text(slide, "中午展示操作稿：真下載、真進度、本機輸出", 0.68, 1.88, 8.5, 0.45, size=22, color=MUTED)
    add_panel(slide, 0.65, 2.65, 5.8, 3.7)
    add_text(slide, "展示主張", 1.0, 2.95, 2.0, 0.36, size=18, color=ACCENT, bold=True)
    add_bullets(
        slide,
        [
            "不是假資料：先連公開 Socrata；若逾時會明確切到備援公開 CSV。",
            "不是假進度：有總大小才顯示 byte 百分比。",
            "展示輸出短路到使用者選擇的資料夾。",
            "產出 payload、manifest、SQLite .db 與 summary JSON。",
        ],
        1.0,
        3.45,
        4.95,
        2.2,
        size=18,
    )
    add_panel(slide, 7.0, 0.95, 5.7, 5.55, fill=RGBColor(8, 12, 18))
    add_text(slide, "GUI 操作入口", 7.35, 1.35, 3.2, 0.45, size=24, color=ACCENT, bold=True)
    add_bullets(
        slide,
        [
            r"scripts\run_showcase_ui.cmd",
            "工具 > 展示模式：下載資料到本機資料夾",
            "工具 > 展示模式：大型 CSV 續傳下載",
            "完成後打開輸出資料夾檢查 .db / manifest / JSON",
        ],
        7.35,
        2.05,
        4.7,
        2.25,
        size=17,
    )
    add_panel(slide, 7.35, 4.8, 4.65, 0.9, fill=PANEL_2)
    add_text(slide, "不放未驗證截圖", 7.65, 5.02, 4.0, 0.3, size=18, color=YELLOW, bold=True)
    add_text(slide, "第一頁只保留可操作入口，避免用錯 IDE 截圖誤導展示。", 7.65, 5.35, 4.0, 0.28, size=13, color=MUTED)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, "本次展示已實作的閉環", 0.65, 0.45, 7, 0.55, size=34, bold=True)
    add_text(slide, "這次只展示已接上 GUI、可重跑、可驗證輸出的功能；仍在開發中的路徑不放進主流程。", 0.68, 1.08, 11.8, 0.35, size=18, color=MUTED)
    steps = [
        ("1", "選擇資料夾", "預設 Downloads，也可選雲端資料夾"),
        ("2", "設定筆數", "由操作者決定小樣本大小"),
        ("3", "真實下載", "公開 Socrata rows API；逾時則備援公開 CSV"),
        ("4", "寫入 manifest", "保留來源與完整性紀錄"),
        ("5", "匯入 .db", "生成 curated_showcase.db"),
    ]
    for index, (number, title, detail) in enumerate(steps):
        x = 0.75 + index * 2.45
        add_panel(slide, x, 2.05, 2.05, 2.35, fill=PANEL_2)
        add_text(slide, number, x + 0.18, 2.23, 0.38, 0.35, size=18, color=ACCENT, bold=True)
        add_text(slide, title, x + 0.18, 2.78, 1.65, 0.36, size=18, bold=True)
        add_text(slide, detail, x + 0.18, 3.25, 1.65, 0.75, size=13, color=MUTED)
        if index < len(steps) - 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 2.05), Inches(3.18), Inches(x + 2.37), Inches(3.18))
            line.line.color.rgb = ACCENT
            line.line.width = Pt(2)
    add_panel(slide, 0.8, 5.15, 11.65, 1.15, fill=RGBColor(18, 28, 39))
    add_text(slide, "現況口徑", 1.1, 5.38, 1.2, 0.3, size=16, color=ACCENT, bold=True)
    add_text(slide, "目前已能從 GUI 觸發公開來源下載，輸出 payload、manifest、SQLite .db 與 summary JSON。", 2.2, 5.35, 9.6, 0.42, size=18)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, "現場操作腳本", 0.65, 0.45, 6, 0.55, size=34, bold=True)
    add_bullets(
        slide,
        [
            r"1. 執行 scripts\run_showcase_ui.cmd，打開 RRKAL Tk 介面。",
            "2. 進入「工具 / 展示模式：下載資料到本機資料夾」。",
            "3. 選擇輸出資料夾；預設使用系統 Downloads。",
            "4. 輸入樣本筆數，例如 100；要壓力測試可調大。",
            "5. 觀察進度視窗：百分比、目前階段、bytes 狀態；若來源逾時會顯示備援切換。",
            "6. 完成後打開輸出資料夾，展示 .db、manifest、summary JSON。",
        ],
        0.9,
        1.35,
        6.2,
        4.8,
        size=19,
    )
    add_panel(slide, 7.35, 1.35, 4.75, 4.6, fill=PANEL_2)
    add_text(slide, "不要現場改程式碼", 7.75, 1.72, 3.6, 0.42, size=22, color=YELLOW, bold=True)
    add_text(slide, "展示模式已包進 GUI。組員只需要看操作流程與輸出結果，不需要接觸 repo、CLI 或測試檔。", 7.75, 2.35, 3.8, 1.2, size=18)
    add_text(slide, "建議展示筆數", 7.75, 4.0, 3.8, 0.35, size=17, color=ACCENT, bold=True)
    add_bullets(slide, ["100：穩定快速", "1,000：較有感", "10,000+：看網路與來源回應"], 7.75, 4.42, 3.8, 1.1, size=15, color=MUTED)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, "進度條不是假的", 0.65, 0.45, 6, 0.55, size=34, bold=True)
    add_text(slide, "GUI 顯示的百分比來自兩種可驗證訊號：流程階段與 HTTP byte 進度。", 0.68, 1.08, 11, 0.35, size=18, color=MUTED)
    add_panel(slide, 0.85, 1.7, 5.7, 4.7, fill=PANEL_2)
    add_text(slide, "流程百分比", 1.15, 2.0, 2.4, 0.35, size=21, color=ACCENT, bold=True)
    for y, pct, label in [(2.75, 5, "準備資料夾"), (3.25, 25, "解析計畫"), (3.75, 35, "開始下載"), (4.25, 82, "下載與匯入完成"), (4.75, 100, "完成")]:
        add_progress(slide, 1.18, y, 4.55, pct, label, color=GREEN if pct >= 82 else ACCENT)
    add_panel(slide, 7.05, 1.7, 5.15, 4.7, fill=PANEL_2)
    add_text(slide, "byte 百分比規則", 7.35, 2.0, 2.8, 0.35, size=21, color=ACCENT, bold=True)
    add_bullets(
        slide,
        [
            "遠端提供 Content-Length：顯示真下載百分比。",
            "遠端沒有總大小：顯示已接收 bytes，不假裝知道總量。",
            "主要公開來源逾時：顯示錯誤並切到備援公開 CSV，不偽裝成功。",
            "完成後用 manifest 與 SQLite 筆數證明成果。",
        ],
        7.35,
        2.65,
        4.15,
        2.25,
        size=18,
    )
    add_text(slide, "展示說法：我們寧可說「總大小未知」，也不要做假的 87%。", 7.35, 5.2, 4.1, 0.55, size=18, color=YELLOW)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, "完成後要展示的輸出", 0.65, 0.45, 6.5, 0.55, size=34, bold=True)
    outputs = [
        ("curated_showcase.db", "本機 SQLite 資料庫，可展示資料真的已匯入。", GREEN),
        ("downloads / manifest.json", "原始 payload 與完整性邊車紀錄。", ACCENT),
        ("showcase_download_summary.json", "機器可讀的展示結果摘要。", YELLOW),
        ("socrata_311.resolved.json", "可稽核的轉接計畫結果。", RED),
    ]
    for index, (title, detail, color) in enumerate(outputs):
        y = 1.45 + index * 1.25
        add_panel(slide, 0.9, y, 11.4, 0.9, fill=PANEL_2)
        add_tag(slide, str(index + 1), 1.15, y + 0.24, color=color)
        add_text(slide, title, 2.0, y + 0.18, 3.8, 0.32, size=20, bold=True)
        add_text(slide, detail, 5.55, y + 0.21, 5.9, 0.32, size=17, color=MUTED)
    add_text(slide, "輸出路徑固定在「你選的資料夾 / RuRuKa Asset Launcher Showcase」。", 0.95, 6.55, 11, 0.35, size=16, color=MUTED)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_text(slide, "已完成 / 未完成邊界", 0.65, 0.45, 6.5, 0.55, size=34, bold=True)
    add_panel(slide, 0.8, 1.35, 5.8, 4.95, fill=PANEL_2)
    add_text(slide, "已完成", 1.15, 1.7, 2.8, 0.4, size=24, color=GREEN, bold=True)
    add_bullets(slide, ["GUI 展示入口已接上。", "樣本筆數可由操作者控制。", "可產出本機 .db、manifest、summary JSON。"], 1.15, 2.35, 4.8, 1.7, size=19)
    add_text(slide, "現場重點：打開 GUI、按展示下載、檢查輸出檔。", 1.15, 4.75, 4.75, 0.6, size=19, color=YELLOW)
    add_panel(slide, 7.1, 1.35, 5.1, 4.95, fill=PANEL_2)
    add_text(slide, "本次不展示", 7.45, 1.7, 3.1, 0.4, size=24, color=ACCENT, bold=True)
    add_bullets(slide, ["全來源無界爬取仍需逐一補安全上限。", "MySQL/PostgreSQL 對接不在本次展示路徑。", "完整錯誤修復 UI 還在硬化。"], 7.45, 2.35, 4.2, 1.7, size=19)
    add_text(slide, "討論重點：目前展示線是否足夠說明進度，以及下一個要補細的功能。", 7.45, 4.75, 4.2, 0.8, size=19)
    return prs


def extracted_texts(path: Path) -> list[str]:
    loaded = Presentation(str(path))
    texts: list[str] = []
    for slide in loaded.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        texts.append(text)
    return texts


def verify_output(path: Path) -> None:
    texts = extracted_texts(path)
    joined = "\n".join(texts)
    markers = ("�", "????", "撅", "嚗")
    bad = [text for text in texts if any(marker in text for marker in markers)]
    if bad:
        raise SystemExit("PPT text verification failed; bad text sample: " + repr(bad[:5]))
    required = ("RuRuKa Asset Launcher", "進度條不是假的", "真下載", "本次展示已實作的閉環")
    missing = [item for item in required if item not in joined]
    if missing:
        raise SystemExit(f"PPT text verification failed; missing text: {missing}")
    print(f"created={path}")
    print(f"slides={len(Presentation(str(path)).slides)}")
    print(f"text_items={len(texts)}")
    print("sample_text=")
    print("\n".join(texts[:12]))


def main() -> int:
    SHOWCASE_DIR.mkdir(parents=True, exist_ok=True)
    deck = build_deck()
    output = OUTPUT
    try:
        deck.save(output)
    except PermissionError:
        # 展示前最常見的真實交付情境是 PPT 正被 PowerPoint 或雲端同步鎖住；
        # 這時不要讓流程整個失敗，改輸出一份帶時間戳的可交付副本。
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output = SHOWCASE_DIR / f"RRKAL_Showcase_Guide.zh-TW.{stamp}.pptx"
        deck.save(output)
        print(f"warning=primary output locked; wrote fallback copy: {output}")
    verify_output(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
